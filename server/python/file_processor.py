import os
import json
import sys
import mimetypes
import re
import shutil
from groq import Groq
from dotenv import load_dotenv

# File processing imports
import pdfplumber
try:
    import pypdfium2 as pdfium
except ImportError:
    pdfium = None

try:
    import pytesseract
except ImportError:
    pytesseract = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from docx import Document
except ImportError:
    Document = None

load_dotenv()

# Set TESSDATA_PREFIX for Tesseract OCR if not already set
if not os.environ.get("TESSDATA_PREFIX"):
    heroku_tessdata = "/app/.apt/usr/share/tesseract-ocr/5/tessdata"
    local_tessdata = "/usr/share/tesseract-ocr/5/tessdata"
    if os.path.exists(heroku_tessdata):
        os.environ["TESSDATA_PREFIX"] = heroku_tessdata
    elif os.path.exists(local_tessdata):
        os.environ["TESSDATA_PREFIX"] = local_tessdata

# 1. Setup Groq Client
client = Groq(
    api_key=os.environ.get("GROQ_API_KEY"),
)

MAX_CHARS_PER_CHUNK = 8000
MAX_CHUNKS = 4


def resolve_tesseract_cmd():
    """Find the tesseract binary in common local, virtualenv, and Heroku apt locations."""
    if pytesseract is None:
        return None

    candidates = [
        os.environ.get("TESSERACT_CMD"),
        os.environ.get("TESSERACT_PATH"),
        shutil.which("tesseract"),
        "/app/.apt/usr/bin/tesseract",
        "/usr/bin/tesseract",
        "/usr/local/bin/tesseract",
    ]

    for candidate in candidates:
        if candidate and os.path.exists(candidate):
            pytesseract.pytesseract.tesseract_cmd = candidate
            return candidate

    return None


def clean_extracted_text(text):
    """Sanitize extracted text to remove corrupted characters while preserving mathematical content."""
    if not text:
        return text
    
    # Remove control characters and other invisible/corrupting characters
    # Keep: newlines, tabs, common punctuation, Unicode symbols
    text = re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f\x7f-\x9f]', '', text)
    
    # Fix common OCR errors: ligatures and common replacements only
    # Don't touch Unicode math symbols - the AI will handle them
    replacements = {
        'ﬁ': 'fi',      # fi ligature
        'ﬂ': 'fl',      # fl ligature
        'ﬀ': 'ff',      # ff ligature
        '–': '-',       # en-dash to hyphen
        '—': '-',       # em-dash to hyphen
        '\u2018': "'",  # left single quote
        '\u2019': "'",  # right single quote
        '\u201c': '"',  # left double quote
        '\u201d': '"',  # right double quote
        '„': '"',
        '‟': '"',
    }
    
    for char, replacement in replacements.items():
        text = text.replace(char, replacement)
    
    # Clean up excessive whitespace but preserve structure
    text = re.sub(r' +', ' ', text)      # Multiple spaces to single
    text = re.sub(r'\n\n+', '\n\n', text)  # Multiple newlines to double
    text = re.sub(r'\t+', '\t', text)    # Multiple tabs to single
    
    # Remove trailing whitespace from lines
    lines = [line.rstrip() for line in text.split('\n')]
    text = '\n'.join(lines)
    
    return text.strip()

# ==================== FILE EXTRACTION FUNCTIONS ====================

def extract_text_from_pdf_ocr(pdf_path):
    """OCR fallback for scanned/image-only PDFs."""
    if pdfium is None or pytesseract is None:
        missing = []
        if pdfium is None:
            missing.append("pypdfium2")
        if pytesseract is None:
            missing.append("pytesseract")
        print(
            f"OCR_UNAVAILABLE: Missing Python OCR dependencies: {', '.join(missing)}.",
            file=sys.stderr,
        )
        return ""

    try:
        tesseract_cmd = resolve_tesseract_cmd()
        if not tesseract_cmd:
            print(
                "OCR_UNAVAILABLE: tesseract binary not found. Install tesseract on the server runtime or set TESSERACT_CMD.",
                file=sys.stderr,
            )
            return ""

        # Verify tesseract binary is accessible.
        _ = pytesseract.get_tesseract_version()
    except Exception:
        print(
            "OCR_UNAVAILABLE: tesseract binary not found. Install tesseract on the server runtime or set TESSERACT_CMD.",
            file=sys.stderr,
        )
        return ""

    ocr_text_parts = []
    try:
        pdf = pdfium.PdfDocument(pdf_path)
        page_count = len(pdf)
        for page_index in range(page_count):
            page = pdf[page_index]
            # Render at higher scale for better OCR quality.
            bitmap = page.render(scale=2.5)
            pil_image = bitmap.to_pil()
            text = pytesseract.image_to_string(pil_image) or ""
            text = text.strip()
            if text:
                ocr_text_parts.append(text)
        ocr_text = "\n".join(ocr_text_parts).strip()
        if ocr_text:
            print(f"OCR extracted {len(ocr_text)} characters from scanned PDF.", file=sys.stderr)
        return ocr_text
    except Exception as e:
        print(f"OCR failed: {e}", file=sys.stderr)
        return ""

def extract_text_from_pdf(pdf_path):
    """Extract text from PDF using pdfplumber."""
    print(f"Extracting text from PDF: {pdf_path}...", file=sys.stderr)
    full_text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = ""

                # Primary extraction.
                basic_text = page.extract_text()
                if basic_text and basic_text.strip():
                    page_text = basic_text

                # Fallback: layout-aware extraction can recover text in some PDFs.
                if not page_text:
                    layout_text = page.extract_text(layout=True)
                    if layout_text and layout_text.strip():
                        page_text = layout_text

                # Fallback: word-level extraction for difficult PDFs.
                if not page_text:
                    words = page.extract_words() or []
                    if words:
                        page_text = " ".join(w.get("text", "") for w in words if w.get("text"))

                # Fallback: character stream extraction helps with some encoded PDFs
                # where higher-level extractors return empty output.
                if not page_text:
                    chars = getattr(page, "chars", None) or []
                    if chars:
                        page_text = "".join(c.get("text", "") for c in chars if c.get("text"))

                if page_text and page_text.strip():
                    full_text += page_text.strip() + "\n"

        full_text = full_text.strip()
        print(f"Extracted {len(full_text)} characters from PDF.", file=sys.stderr)

        if len(full_text) == 0:
            # OCR fallback for scanned/image-only PDFs.
            ocr_text = extract_text_from_pdf_ocr(pdf_path)
            if ocr_text:
                return ocr_text

            print(
                "NO_EXTRACTABLE_TEXT: PDF appears scanned/image-only or has no selectable text.",
                file=sys.stderr,
            )
            return None

        return full_text
    except Exception as e:
        print(f"Error reading PDF: {e}", file=sys.stderr)
        return None


def extract_text_from_pptx(pptx_path):
    """Extract text from PowerPoint presentation."""
    print(f"Extracting text from PPTX: {pptx_path}...", file=sys.stderr)
    if Presentation is None:
        raise RuntimeError("python-pptx not installed. Install with: pip install python-pptx")
    
    full_text = ""
    try:
        prs = Presentation(pptx_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            full_text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        full_text += shape.text + "\n"
        
        print(f"Extracted {len(full_text)} characters from PPTX.", file=sys.stderr)
        return full_text
    except Exception as e:
        print(f"Error reading PPTX: {e}", file=sys.stderr)
        return None


def extract_text_from_docx(docx_path):
    """Extract text from Word document."""
    print(f"Extracting text from DOCX: {docx_path}...", file=sys.stderr)
    if Document is None:
        raise RuntimeError("python-docx not installed. Install with: pip install python-docx")
    
    full_text = ""
    try:
        doc = Document(docx_path)
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                full_text += paragraph.text + "\n"
        
        print(f"Extracted {len(full_text)} characters from DOCX.", file=sys.stderr)
        return full_text
    except Exception as e:
        print(f"Error reading DOCX: {e}", file=sys.stderr)
        return None


def extract_text_from_plain_text(file_path):
    """Extract text from plain text files (TXT, MD, MARKDOWN, CSV, JSON, RTF)."""
    print(f"Reading text file: {file_path}...", file=sys.stderr)
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            full_text = f.read()
        print(f"Extracted {len(full_text)} characters.", file=sys.stderr)
        return full_text
    except Exception as e:
        print(f"Error reading text file: {e}", file=sys.stderr)
        return None


def extract_text_from_file(file_path):
    """
    Detect file type and extract text accordingly.
    Supports: PDF, PPTX, DOCX, TXT, MD, MARKDOWN, CSV, JSON, RTF
    """
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    print(f"Processing file: {file_path} (type: {ext})", file=sys.stderr)
    
    if ext == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif ext == ".pptx":
        text = extract_text_from_pptx(file_path)
    elif ext == ".docx":
        text = extract_text_from_docx(file_path)
    elif ext in [".txt", ".md", ".markdown", ".csv", ".json", ".rtf"]:
        text = extract_text_from_plain_text(file_path)
    else:
        # Try to read as plain text for unknown formats
        print(f"Unknown format {ext}, attempting to read as text...", file=sys.stderr)
        text = extract_text_from_plain_text(file_path)
    
    # Clean the extracted text to sanitize corrupted characters
    if text:
        text = clean_extracted_text(text)
    
    return text


# ==================== TEXT CHUNKING & DEDUPLICATION ====================

def chunk_text(text, max_chars=MAX_CHARS_PER_CHUNK, max_chunks=MAX_CHUNKS):
    """Split long text into chunks to reduce truncation and improve coverage."""
    if len(text) <= max_chars:
        return [text]

    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks = []
    current = ""

    for para in paragraphs:
        candidate = f"{current}\n\n{para}".strip() if current else para
        if len(candidate) <= max_chars:
            current = candidate
            continue

        if current:
            chunks.append(current)
            if len(chunks) >= max_chunks:
                break

        # Handle very large single paragraphs.
        if len(para) > max_chars:
            start = 0
            while start < len(para) and len(chunks) < max_chunks:
                chunks.append(para[start:start + max_chars])
                start += max_chars
            current = ""
        else:
            current = para

    if current and len(chunks) < max_chunks:
        chunks.append(current)

    return chunks[:max_chunks]


def dedupe_study_set(study_set):
    """Deduplicate cards/questions while preserving order."""
    flashcards = []
    seen_cards = set()
    for card in study_set.get("flashcards", []):
        front = str(card.get("front", "")).strip()
        back = str(card.get("back", "")).strip()
        if not front or not back:
            continue
        key = (front, back)
        if key in seen_cards:
            continue
        seen_cards.add(key)
        flashcards.append({"front": front, "back": back})

    quiz = []
    seen_quiz = set()
    for item in study_set.get("quiz", []):
        question = str(item.get("question", "")).strip()
        options = item.get("options", [])
        correct_answer = str(item.get("correct_answer", "")).strip()
        if not question or not isinstance(options, list) or len(options) < 2 or not correct_answer:
            continue
        normalized_options = [str(o).strip() for o in options if str(o).strip()]
        key = (question, tuple(normalized_options), correct_answer)
        if key in seen_quiz:
            continue
        seen_quiz.add(key)
        quiz.append(
            {
                "question": question,
                "options": normalized_options,
                "correct_answer": correct_answer,
            }
        )

    return {"flashcards": flashcards, "quiz": quiz}


def _ensure_math_wrapping(value):
    """Best-effort guard so math-like fragments are wrapped for KaTeX rendering."""
    text = str(value or "")
    if not text:
        return text

    if "$" in text or "\\(" in text or "\\[" in text:
        return text

    # Convert common delimiters first.
    text = re.sub(r"\\\((.*?)\\\)", lambda m: f"${m.group(1).strip()}$", text)
    text = re.sub(r"\\\[(.*?)\\\]", lambda m: f"$${m.group(1).strip()}$$", text)

    # Wrap math-like token clusters (variables with subscripts/superscripts, latex commands).
    pattern = re.compile(
        r"([A-Za-z][A-Za-z0-9]*(?:_[A-Za-z0-9]+|\^[A-Za-z0-9{}()+\-]+)+(?:\([^)]+\))?"
        r"|\\(?:int|frac|sum|prod|lim|to|infty|cdot|times|le|ge|ne|sqrt|alpha|beta|gamma|delta|theta|lambda|pi|sigma|omega)\b(?:\s*\{[^}]+\})?)"
    )

    def repl(match):
        seg = match.group(1)
        return f"${seg}$"

    return pattern.sub(repl, text)


def ensure_math_in_study_set(study_set):
    for card in study_set.get("flashcards", []):
        card["front"] = _ensure_math_wrapping(card.get("front", ""))
        card["back"] = _ensure_math_wrapping(card.get("back", ""))

    for item in study_set.get("quiz", []):
        item["question"] = _ensure_math_wrapping(item.get("question", ""))
        item["correct_answer"] = _ensure_math_wrapping(item.get("correct_answer", ""))
        options = item.get("options", [])
        if isinstance(options, list):
            item["options"] = [_ensure_math_wrapping(opt) for opt in options]

    return study_set


# ==================== GROQ GENERATION ====================

def _generate_study_material_for_chunk(context_text, chunk_index, total_chunks):
    print(
        f"Generating chunk {chunk_index + 1}/{total_chunks} via Groq Structured Outputs...",
        file=sys.stderr,
    )

    # Define the JSON Schema
    study_material_schema = {
        "type": "object",
        "properties": {
            "flashcards": {
                "type": "array",
                "minItems": 1,
                "items": {
                    "type": "object",
                    "properties": {
                        "front": {"type": "string"},
                        "back": {"type": "string"}
                    },
                    "required": ["front", "back"],
                    "additionalProperties": False
                }
            },
            "quiz": {
                "type": "array",
                "minItems": 0,
                "items": {
                    "type": "object",
                    "properties": {
                        "question": {"type": "string"},
                        "options": {
                            "type": "array",
                            "items": {"type": "string"},
                            "minItems": 4,
                            "maxItems": 4
                        },
                        "correct_answer": {"type": "string"}
                    },
                    "required": ["question", "options", "correct_answer"],
                    "additionalProperties": False
                }
            }
        },
        "required": ["flashcards", "quiz"],
        "additionalProperties": False
    }

    # Fallback schema for cases where the model repeatedly omits quiz output.
    flashcards_only_schema = {
        "type": "object",
        "properties": {
            "flashcards": {
                "type": "array",
                "minItems": 1,
                "items": {
                    "type": "object",
                    "properties": {
                        "front": {"type": "string"},
                        "back": {"type": "string"}
                    },
                    "required": ["front", "back"],
                    "additionalProperties": False
                }
            }
        },
        "required": ["flashcards"],
        "additionalProperties": False
    }

    try:
        prompt = (
            "Create concise, high-quality study materials from these notes. "
            "Cover distinct concepts and avoid duplicates. "
            "CRITICAL MATH FORMATTING RULES:\n"
            "1. Convert ALL Unicode math symbols (α β γ δ ε ζ η θ λ μ ξ π ρ σ τ υ φ ψ ω Γ Δ Λ Π Σ Ω) to LaTeX commands (\\alpha \\beta, etc.)\n"
            "2. Wrap ALL mathematical notation with $ $ for inline or $$ $$ for display math.\n"
            "3. Use LaTeX commands: \\int \\sum \\prod \\frac{a}{b} \\sqrt{x} \\times \\div \\le \\ge \\ne \\infty \\pm etc.\n"
            "4. Use subscripts_{x} and superscripts^{2} with proper braces.\n"
            "5. Never output raw unicode math symbols - convert to LaTeX equivalents.\n"
            "6. Every variable, equation, and symbol must be properly delimited.\n\n"
            "For flashcards: use clear term/question on front and precise explanation on back. "
            "For quiz: include 4 plausible options and ensure correct_answer exactly matches one option.\n\n"
            f"Chunk {chunk_index + 1} of {total_chunks}.\n"
            f"Lecture Notes:\n{context_text}"
        )

        last_err = None
        # First attempt full schema (flashcards + quiz).
        for _ in range(2):
            try:
                chat_completion = client.chat.completions.create(
                    messages=[
                        {
                            "role": "system",
                            "content": "You are a helpful study assistant. Return only schema-compliant JSON. "
                            "ALL math must use LaTeX formatting with $...$ delimiters. "
                            "Convert Unicode math symbols to LaTeX commands immediately.",
                        },
                        {"role": "user", "content": prompt},
                    ],
                    model="openai/gpt-oss-120b",
                    temperature=0.0,
                    response_format={
                        "type": "json_schema",
                        "json_schema": {
                            "name": "study_set",
                            "strict": True,
                            "schema": study_material_schema,
                        },
                    },
                )
                parsed = json.loads(chat_completion.choices[0].message.content)
                return ensure_math_in_study_set(parsed)
            except Exception as inner:
                last_err = inner

        # Second attempt: flashcards-only schema, then normalize quiz to an empty array.
        for _ in range(2):
            try:
                chat_completion = client.chat.completions.create(
                    messages=[
                        {
                            "role": "system",
                            "content": "You are a helpful study assistant. Return only schema-compliant JSON. "
                            "Math must always be wrapped in $...$ or $$...$$.",
                        },
                        {
                            "role": "user",
                            "content": (
                                prompt
                                + "\n\nIf quiz generation is uncertain, return only flashcards."
                            ),
                        },
                    ],
                    model="openai/gpt-oss-120b",
                    temperature=0.0,
                    response_format={
                        "type": "json_schema",
                        "json_schema": {
                            "name": "study_set_flashcards_only",
                            "strict": True,
                            "schema": flashcards_only_schema,
                        },
                    },
                )
                parsed = json.loads(chat_completion.choices[0].message.content)
                parsed["quiz"] = []
                return ensure_math_in_study_set(parsed)
            except Exception as inner:
                last_err = inner

        raise last_err if last_err else RuntimeError("Unknown Groq generation error")

    except Exception as e:
        print(f"Error calling Groq API: {e}", file=sys.stderr)
        return None


def generate_study_material(context_text):
    chunks = chunk_text(context_text)
    print(f"Processing {len(chunks)} chunk(s).", file=sys.stderr)

    merged = {"flashcards": [], "quiz": []}
    successful_chunks = 0
    for idx, chunk in enumerate(chunks):
        result = _generate_study_material_for_chunk(chunk, idx, len(chunks))
        if not result:
            continue
        successful_chunks += 1
        merged["flashcards"].extend(result.get("flashcards", []))
        merged["quiz"].extend(result.get("quiz", []))

    merged = dedupe_study_set(merged)
    if successful_chunks == 0 or (len(merged["flashcards"]) == 0 and len(merged["quiz"]) == 0):
        print(
            "NO_STUDY_MATERIALS: Model generation returned no usable flashcards or quiz.",
            file=sys.stderr,
        )
        return None

    print(
        f"Generated {len(merged['flashcards'])} flashcards and {len(merged['quiz'])} quiz items.",
        file=sys.stderr,
    )
    return json.dumps(merged)


# ==================== MAIN PIPELINE ====================

def run_pipeline(file_path):
    """
    Extract text from any supported file format and generate study materials.
    Returns JSON string or None on failure.
    When run as CLI: accepts file_path as arg, prints JSON to stdout (errors to stderr).
    """
    raw_text = extract_text_from_file(file_path)
    if not raw_text:
        return None
    return generate_study_material(raw_text)


def main():
    if len(sys.argv) < 2:
        print("Usage: python file_processor.py <path_to_file>", file=sys.stderr)
        sys.exit(1)

    file_path = sys.argv[1]
    json_output = run_pipeline(file_path)

    if json_output:
        # Print JSON to stdout for Node.js to capture
        print(json_output)
    else:
        # Exit with error code
        sys.exit(1)


if __name__ == "__main__":
    main()
